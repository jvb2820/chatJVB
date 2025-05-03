import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import CohereEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain_cohere import ChatCohere
import os
import base64
import io
from PIL import Image
import fitz  # PyMuPDF for PDF to image conversion
import docx  # python-docx for PDF to Word conversion
import re
from pptx import Presentation  # python-pptx for PowerPoint conversion
from pptx.util import Inches, Pt

# Cohere API Key
COHERE_API_KEY = "aFR2rly7rpnQoOJ4Xxo1n6dAz4whPkemrnvztoA7"

# Set logo paths (adjust as needed)
CHATBOT_LOGO_PATH = os.path.join(os.path.dirname(__file__), "logo.jpg")  # Chatbot logo
USER_LOGO_PATH = os.path.join(os.path.dirname(__file__), "user_logo.svg")  # User logo

# Pre-defined responses for specific questions
PRE_DEFINED_RESPONSES = {
    "who is jvb": "JVB is Jeuz Vinci Bas, my creator. He is a talented programmer and software developer. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/. He specializes in AI applications and software development.",
    
    "what is jvb": "JVB stands for Jeuz Vinci Bas, the developer who created me. He's a programmer with expertise in AI and software development. Feel free to connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.",
    
    "tell me about jvb": "Jeuz Vinci Bas (JVB) is a programmer and developer who created me. He works on various software projects including AI applications like this chatbot. You can learn more about his work or connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.",
    
    "jvb": "JVB refers to Jeuz Vinci Bas, my creator and a software developer. He built this chatbot as part of his work in AI development. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/ if you'd like to learn more about his projects.",
    
    "who created you": "I was created by Jeuz Vinci Bas (JVB), a programmer and AI developer. This chatbot is one of his projects combining AI conversation capabilities with practical tools like PDF conversion.",
    
    # New entries for variations of the name
    "who is jeuz": "Jeuz Vinci Bas is my creator. He is a talented programmer and software developer with expertise in AI applications. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.",
    
    "who is jeuz vinci bas": "Jeuz Vinci Bas is my creator and developer. He specializes in AI applications and software development, including this chatbot. You can learn more about his work or connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.",
    
    "tell me about jeuz": "Jeuz Vinci Bas is the developer who created me. He works on various software projects with a focus on AI applications. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/ to learn more about his work.",
    
    "tell me about jeuz vinci bas": "Jeuz Vinci Bas is a programmer and AI developer who created this chatbot. He combines AI conversation capabilities with practical tools like PDF conversion in his projects. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.",
    
    "jeuz": "Jeuz refers to Jeuz Vinci Bas, my creator and a software developer. He built this chatbot as part of his work in AI development. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.",
    
    "jeuz vinci bas": "Jeuz Vinci Bas is my creator and a talented software developer specializing in AI applications. This chatbot is one of his projects. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/."
}

# Function to check for pre-defined responses with improved matching
def get_predefined_response(question):
    # Convert question to lowercase for case-insensitive matching
    question_lower = question.lower().strip()
    
    # Direct match check
    if question_lower in PRE_DEFINED_RESPONSES:
        return PRE_DEFINED_RESPONSES[question_lower]
    
    # Create a normalized version of the question (remove punctuation)
    normalized_question = re.sub(r'[^\w\s]', '', question_lower)
    normalized_question_words = set(normalized_question.split())
    
    # Check for 'jeuz' variations explicitly
    if ('who' in normalized_question_words and 'jeuz' in normalized_question_words) or \
       ('tell' in normalized_question_words and 'about' in normalized_question_words and 'jeuz' in normalized_question_words) or \
       ('jeuz' in normalized_question_words and len(normalized_question_words) <= 3):
        return PRE_DEFINED_RESPONSES.get("who is jeuz", 
               "Jeuz Vinci Bas is my creator. He is a talented programmer and software developer with expertise in AI applications. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/.")
    
    # Check for 'jvb' variations similarly
    if ('who' in normalized_question_words and 'jvb' in normalized_question_words) or \
       ('tell' in normalized_question_words and 'about' in normalized_question_words and 'jvb' in normalized_question_words) or \
       ('jvb' in normalized_question_words and len(normalized_question_words) <= 3):
        return PRE_DEFINED_RESPONSES.get("who is jvb",
               "JVB is Jeuz Vinci Bas, my creator. He is a talented programmer and software developer. You can connect with him on LinkedIn at https://www.linkedin.com/in/jeuz-vinci-bas-b51639341/. He specializes in AI applications and software development.")
    
    # Check if any key is contained within the question
    for key, response in PRE_DEFINED_RESPONSES.items():
        # For specific question patterns, check if all words in the key appear in the question
        key_words = set(key.split())
        
        # If all key words are in the question
        if key_words.issubset(normalized_question_words):
            return response
            
    # No match found
    return None

# Custom CSS for styling the chat messages
st.markdown("""
<style>
    .user-message {
        background-color: #123C5A;
        padding: 10px;
        border-radius: 10px;
        margin-bottom: 10px;
        display: flex;
        align-items: flex-start;

    }
    .bot-message {
        background-color: #0A1D37;
        padding: 10px;;
        padding: 10px;
        border-radius: 10px;
        margin-bottom: 10px;
        display: flex;
        align-items: flex-start;
    }
    .message-container {
        margin-bottom: 15px;
    }
    .logo-container {
        margin-right: 10px;
        flex-shrink: 0;
    }
    .logo-image {
        width: 30px;
        height: 30px;
        border-radius: 50%;
        object-fit: cover;
    }
    .message-content {
        flex-grow: 1;
    }
</style>
""", unsafe_allow_html=True)

# App header
st.header("🚀 ChatJVB")

# Initialize session state for storing conversation history
if "conversation_history" not in st.session_state:
    st.session_state.conversation_history = []

# Create a container for chat history display
chat_container = st.container()

# PDF to Word conversion function
def convert_pdf_to_word(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    doc = docx.Document()
    
    for page_num in range(len(pdf_reader.pages)):
        text = pdf_reader.pages[page_num].extract_text()
        doc.add_paragraph(text)
    
    # Save to a BytesIO object
    docx_bytes = io.BytesIO()
    doc.save(docx_bytes)
    docx_bytes.seek(0)
    
    return docx_bytes

# PDF to Image conversion function
def convert_pdf_to_images(pdf_file):
    pdf_bytes = pdf_file.read()
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    images = []
    
    for page_num in range(len(pdf_document)):
        page = pdf_document.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better quality
        img_bytes = pix.pil_tobytes(format="JPEG")
        img = Image.open(io.BytesIO(img_bytes))
        images.append(img)
    
    return images

# PDF to PowerPoint conversion function
def convert_pdf_to_powerpoint(pdf_file):
    # Create a new PowerPoint presentation
    prs = Presentation()
    
    # Extract text and images from PDF
    pdf_bytes = pdf_file.read()
    pdf_document = fitz.open(stream=pdf_bytes, filetype="pdf")
    
    # For each page in the PDF
    for page_num in range(len(pdf_document)):
        # Add a slide
        slide_layout = prs.slide_layouts[5]  # Using blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Get page content
        page = pdf_document.load_page(page_num)
        text = page.get_text()
        
        # Add a title to the slide (using first line as title or page number)
        title_text = text.split('\n', 1)[0] if text else f"Slide {page_num + 1}"
        if slide.shapes.title:
            slide.shapes.title.text = title_text
        
        # Add text content to the slide
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Add main content (skip first line as it's used for title)
        content = text.split('\n', 1)[1] if '\n' in text else ""
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(12)
        
        # Alternative image extraction method using pixmap
        try:
            # Render the page to an image
            pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5))  # Use a good resolution
            img_bytes = pix.tobytes("png")  # Convert to PNG format
            
            # Save image to a BytesIO object for adding to slide
            image_stream = io.BytesIO(img_bytes)
            image_stream.seek(0)
            
            # Add as background image to maintain layout
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Add the page image as a full slide background
            slide.shapes.add_picture(image_stream, 0, 0, width=slide_width, height=slide_height)
            
            # If we added the image as background, add a semi-transparent overlay for text
            # to ensure text is visible over the image
            left = 0
            top = 0
            width = slide_width
            height = slide_height
            
            # Add a semi-transparent white shape behind text for readability
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, width, height
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            shape.fill.transparency = 0.7  # 70% transparent
            shape.line.color.rgb = RGBColor(255, 255, 255)
            shape.line.transparency = 1.0  # Fully transparent border
            
            # Send the background shape to back
            shape.z_order = 0
            
        except Exception as e:
            print(f"Error rendering page {page_num + 1} as image: {e}")
    
    # Save presentation to BytesIO object
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    
    return pptx_bytes

# Sidebar for PDF conversion tools
with st.sidebar: 
    st.title("🔄 PDF Converter")
    
    pdf_file = st.file_uploader("Upload a PDF file to convert", type="pdf")
    
    if pdf_file is not None:
        # Conversion options
        conversion_type = st.radio(
            "Choose conversion type:",
            ("PDF to Word", "PDF to PowerPoint", "PDF to Images")
        )
        
        if st.button("Convert"):
            with st.spinner("Converting PDF..."):
                if conversion_type == "PDF to Word":
                    docx_bytes = convert_pdf_to_word(pdf_file)
                    st.download_button(
                        label="Download Word Document",
                        data=docx_bytes,
                        file_name=f"{pdf_file.name.split('.pdf')[0]}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                    )
                    st.success("Conversion complete! Click the download button above.")
                
                elif conversion_type == "PDF to PowerPoint":
                    pptx_bytes = convert_pdf_to_powerpoint(pdf_file)
                    st.download_button(
                        label="Download PowerPoint Presentation",
                        data=pptx_bytes,
                        file_name=f"{pdf_file.name.split('.pdf')[0]}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    st.success("Conversion complete! Click the download button above.")
                
                elif conversion_type == "PDF to Images":
                    images = convert_pdf_to_images(pdf_file)
                    st.success(f"Converted {len(images)} pages to images.")
                    
                    # Display and provide download for each image
                    for i, img in enumerate(images):
                        col1, col2 = st.columns([3, 1])
                        with col1:
                            # FIXED: Changed use_column_width to use_container_width
                            st.image(img, caption=f"Page {i+1}", use_container_width=True)
                        
                        with col2:
                            # Convert image to bytes for download
                            img_byte_arr = io.BytesIO()
                            img.save(img_byte_arr, format='PNG')
                            img_byte_arr.seek(0)
                            
                            st.download_button(
                                label="Download",
                                data=img_byte_arr,
                                file_name=f"page_{i+1}.png",
                                mime="image/png"
                            )
    
    # Add clear chat button in sidebar
    if st.button("Clear Chat History"):
        st.session_state.conversation_history = []
        st.rerun()
    
    # Add information about functionality
    st.markdown("---")
    st.write("**How to use:**")
    st.write("1. Chat with the assistant anytime")
    st.write("2. Convert PDFs to Word documents, PowerPoint presentations, or images")

# Define the LLM using Cohere
llm = ChatCohere(
    cohere_api_key=COHERE_API_KEY,
    temperature=0.2,
    max_tokens=1000,
    model="command"
)

# Use form to capture user input with auto-clearing functionality
with st.form(key="question_form", clear_on_submit=True):
    user_question = st.text_input("Type your question here")
    submit_button = st.form_submit_button("Send")

# Handle the user question
if submit_button and user_question:
    # Add user question to history
    st.session_state.conversation_history.append({"role": "user", "content": user_question})
    
    # Check for pre-defined responses first
    predefined_response = get_predefined_response(user_question)
    
    if predefined_response:
        # Use the pre-defined response
        response = predefined_response
    else:
        # Generate response using Cohere if no pre-defined response exists
        response = llm.invoke(user_question).content
    
    # Add response to history
    st.session_state.conversation_history.append({"role": "assistant", "content": response})
    
    # Force a rerun to update the display immediately and clear the input
    st.rerun()

# Load the logos (only once with better error handling)
if "chatbot_logo_base64" not in st.session_state or "user_logo_base64" not in st.session_state:
    try:
        # Load chatbot logo
        with open(CHATBOT_LOGO_PATH, "rb") as f:
            chatbot_logo_bytes = f.read()
            st.session_state.chatbot_logo_base64 = base64.b64encode(chatbot_logo_bytes).decode()

        # Load user logo
        with open(USER_LOGO_PATH, "rb") as f:
            user_logo_bytes = f.read()
            st.session_state.user_logo_base64 = base64.b64encode(user_logo_bytes).decode()
    except Exception as e:
        st.sidebar.warning(f"Could not load logos: {e}")
        # Set fallback empty strings for the logos
        st.session_state.chatbot_logo_base64 = ""
        st.session_state.user_logo_base64 = ""

# Display conversation history with custom styling
with chat_container:
    for message in st.session_state.conversation_history:
        if message["role"] == "user":
            st.markdown(f"""
            <div class="message-container">
                <div class="user-message">
                    <div class="logo-container">
                        <img src="data:image/svg+xml;base64,{st.session_state.get('user_logo_base64', '')}" class="logo-image" alt="User Logo">
                    </div>
                     {message['content']}
                </div>
            </div>
            """, unsafe_allow_html=True)
        else:
            # Display chatbot logo beside the text
            st.markdown(f"""
            <div class="message-container">
                <div class="bot-message">
                    <div class="logo-container">
                        <img src="data:image/png;base64,{st.session_state.get('chatbot_logo_base64', '')}" class="logo-image" alt="Chatbot Logo">
                    </div>
                    <div class="message-content">
                        {message['content']}
                    </div>
                </div>
            </div>
            """, unsafe_allow_html=True)